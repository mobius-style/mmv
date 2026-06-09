#!/usr/bin/env python3
"""
custom_rag_adapter.py — MOBIUS MMV Phase C: Box A Retrieval Adapter
src/adapters/custom_rag_adapter.py

User document Custom RAG adapter.

Design principles (phase_c_spec_v2_2.docx §1, §4):
  Principle F  Index file scheme: custom_index.faiss is the source of truth
  Principle G  CPU/GPU separation: embedding uses MMV_EMBEDDING_DEVICE
               (auto/cpu/cuda), defaulting to CUDA when available.
  Atomic swap: os.replace() for custom_index_new.faiss -> custom_index.faiss
  FileWatcher: watchdog detects corpus changes -> asynchronous rebuild

Supported formats:
  Text         .txt .md .rst .yaml .yml
  PDF          .pdf (text extraction + Tesseract OCR for scanned pages)
  Image        .png .jpg .jpeg .tiff .bmp .webp (Tesseract OCR: eng+jpn)
  Word         .docx .doc
  Excel        .xlsx .xls .csv .tsv
  PowerPoint   .pptx .ppt
  HTML         .html .htm
  Email        .eml .msg (Outlook)
  Data         .json .jsonl
  Rich text    .rtf
  LibreOffice  .odt
  Archive      .zip .7z .rar (recursively extracts internal files)

Technical specifications (phase_c_spec_v2_2.docx §4):
  Index           : data/box_a/custom_index.faiss (IndexFlatIP)
  Chunks          : data/box_a/custom_chunks.jsonl
  Manifest        : data/box_a/index_manifest.json
  Embedding model : intfloat/multilingual-e5-large (device auto-select)
  Chunk size      : 512 tokens / 64 token overlap (per §5.2)
  top_k           : default 5

Author : Taiko Toeda / MOBIUS LLC
License: AGPL-3.0-or-later
Spec   : phase_c_spec_v2_2.docx §4, §5, §7, §11
"""

from __future__ import annotations

import email
import hashlib
import json
import logging
import os
import re
import shutil
import subprocess
import tempfile
import threading
import time
from dataclasses import dataclass
from datetime import datetime, timezone
from pathlib import Path
from typing import List, Optional

import faiss
import numpy as np
from sentence_transformers import SentenceTransformer

logger = logging.getLogger(__name__)

# -- Constants ---------------------------------------------------------------
MODEL_NAME      = "intfloat/multilingual-e5-large"
CHUNK_TOKENS    = 512    # §5.2: 256-512 tokens per chunk
OVERLAP_TOKENS  = 64     # §5.2: 64-token overlap
CHUNK_CHARS     = CHUNK_TOKENS  * 4   # 2048 (~4 chars/token)
OVERLAP_CHARS   = OVERLAP_TOKENS * 4  # 256
MIN_CHUNK_CHARS = 80
EMBED_BATCH     = 32
DEFAULT_TOP_K   = 5
MAX_TOP_K       = 20
OCR_MIN_TEXT    = 50     # Minimum char count to determine if PDF text extraction needs OCR


def _select_embedding_device() -> str:
    requested = os.environ.get("MMV_EMBEDDING_DEVICE", "auto").strip().lower()
    if requested == "cpu":
        return "cpu"
    if requested.startswith("cuda"):
        try:
            import torch
            if torch.cuda.is_available():
                return requested
            logger.warning(
                "MMV_EMBEDDING_DEVICE=%s requested but CUDA is unavailable; using cpu.",
                requested,
            )
        except Exception as exc:
            logger.warning("Could not inspect CUDA availability: %s; using cpu.", exc)
        return "cpu"
    if requested != "auto":
        logger.warning(
            "Unknown MMV_EMBEDDING_DEVICE=%s; expected auto/cpu/cuda. Using auto.",
            requested,
        )
    try:
        import torch
        return "cuda" if torch.cuda.is_available() else "cpu"
    except Exception:
        return "cpu"

# -- Type definitions (per 08_api_types.docx §9) ----------------------------

@dataclass
class Source:
    source_type:     str
    label:           str
    uri:             str
    chunk_index:     Optional[int] = None
    retrieved_at:    str           = ""
    relevance_score: Optional[float] = None

    def __post_init__(self):
        if not self.retrieved_at:
            self.retrieved_at = datetime.now(timezone.utc).isoformat()


@dataclass
class RetrievalResult:
    sources:   List[Source]
    outcome:   str
    synthesis: str


class RetrievalAdapter:
    def retrieve(self, query: str, top_k: int = DEFAULT_TOP_K) -> RetrievalResult:
        raise NotImplementedError
    def is_available(self) -> bool:
        raise NotImplementedError
    def get_sufficiency_score(self, result: RetrievalResult) -> float:
        raise NotImplementedError


# -- Text extraction ---------------------------------------------------------

class _TextExtractor:
    """
    Extracts plain text from various file formats.
    Returns empty string on extraction failure (does not propagate exceptions).
    """

    # Formats that are recursively extracted from archives
    ARCHIVE_EXTENSIONS = {".zip", ".7z", ".rar"}

    # Image extensions for OCR
    IMAGE_EXTENSIONS = {".png", ".jpg", ".jpeg", ".tiff", ".tif", ".bmp", ".webp"}

    @staticmethod
    def extract(path: Path, _depth: int = 0) -> str:
        """
        Extract text from a path.
        Archives are recursively extracted (max 2 levels deep).
        """
        if _depth > 2:
            return ""
        ext = path.suffix.lower()
        try:
            if ext in (".txt", ".md", ".rst"):
                return _TextExtractor._text(path)
            elif ext in (".yaml", ".yml"):
                return _TextExtractor._yaml(path)
            elif ext == ".pdf":
                return _TextExtractor._pdf(path)
            elif ext in _TextExtractor.IMAGE_EXTENSIONS:
                return _TextExtractor._ocr(path)
            elif ext == ".docx":
                return _TextExtractor._docx(path)
            elif ext == ".doc":
                return _TextExtractor._pandoc(path, "plain")
            elif ext in (".xlsx", ".xls"):
                return _TextExtractor._excel(path)
            elif ext in (".csv", ".tsv"):
                return _TextExtractor._csv(path, ext)
            elif ext == ".pptx":
                return _TextExtractor._pptx(path)
            elif ext == ".ppt":
                return _TextExtractor._pandoc(path, "plain")
            elif ext in (".html", ".htm"):
                return _TextExtractor._html(path)
            elif ext == ".eml":
                return _TextExtractor._eml(path)
            elif ext == ".msg":
                return _TextExtractor._msg(path)
            elif ext == ".json":
                return _TextExtractor._json(path)
            elif ext == ".jsonl":
                return _TextExtractor._jsonl(path)
            elif ext == ".rtf":
                return _TextExtractor._rtf(path)
            elif ext == ".odt":
                return _TextExtractor._odt(path)
            elif ext in _TextExtractor.ARCHIVE_EXTENSIONS:
                return _TextExtractor._archive(path, _depth)
            else:
                logger.debug(f"[Extractor] unsupported: {path.name}")
                return ""
        except Exception as e:
            logger.warning(f"[Extractor] failed {path.name}: {e}")
            return ""

    # -- Individual extraction methods ----------------------------------------

    @staticmethod
    def _text(path: Path) -> str:
        for enc in ("utf-8", "utf-8-sig", "cp932", "latin-1"):
            try:
                return path.read_text(encoding=enc)
            except (UnicodeDecodeError, LookupError):
                continue
        return path.read_text(encoding="utf-8", errors="replace")

    @staticmethod
    def _yaml(path: Path) -> str:
        import yaml
        with open(path, encoding="utf-8", errors="replace") as f:
            obj = yaml.safe_load(f)
        return json.dumps(obj, ensure_ascii=False, indent=2) if obj else ""

    @staticmethod
    def _pdf(path: Path) -> str:
        """Text extraction -> falls back to Tesseract OCR if insufficient"""
        from pypdf import PdfReader
        reader = PdfReader(str(path))
        parts = []
        for page in reader.pages:
            t = page.extract_text()
            if t:
                parts.append(t)
        text = "\n".join(parts).strip()
        if len(text) >= OCR_MIN_TEXT:
            return text
        # Scanned PDF -> OCR
        logger.info(f"  [Extractor] PDF text too short ({len(text)} chars). Trying OCR: {path.name}")
        return _TextExtractor._pdf_ocr(path)

    @staticmethod
    def _pdf_ocr(path: Path) -> str:
        """Convert PDF to images and run Tesseract OCR"""
        try:
            import pytesseract
            from PIL import Image
            # Use pdf2image if available
            try:
                from pdf2image import convert_from_path
                images = convert_from_path(str(path), dpi=200)
                parts = []
                for img in images:
                    parts.append(pytesseract.image_to_string(img, lang="eng+jpn"))
                return "\n".join(parts)
            except ImportError:
                # pdf2image not available -> PyMuPDF fallback
                try:
                    import fitz  # PyMuPDF
                    doc = fitz.open(str(path))
                    parts = []
                    for page in doc:
                        pix = page.get_pixmap(dpi=200)
                        img_data = pix.tobytes("png")
                        from io import BytesIO
                        img = Image.open(BytesIO(img_data))
                        parts.append(pytesseract.image_to_string(img, lang="eng+jpn"))
                    return "\n".join(parts)
                except ImportError:
                    logger.warning(f"  [Extractor] pdf2image and PyMuPDF not available. OCR skipped.")
                    return ""
        except Exception as e:
            logger.warning(f"  [Extractor] PDF OCR failed: {e}")
            return ""

    @staticmethod
    def _ocr(path: Path) -> str:
        """Run Tesseract OCR on image files (eng+jpn)"""
        import pytesseract
        from PIL import Image
        img = Image.open(str(path))
        return pytesseract.image_to_string(img, lang="eng+jpn")

    @staticmethod
    def _docx(path: Path) -> str:
        from docx import Document
        doc = Document(str(path))
        parts = []
        for para in doc.paragraphs:
            if para.text.strip():
                parts.append(para.text)
        for table in doc.tables:
            for row in table.rows:
                row_text = " | ".join(
                    cell.text.strip() for cell in row.cells if cell.text.strip()
                )
                if row_text:
                    parts.append(row_text)
        return "\n".join(parts)

    @staticmethod
    def _excel(path: Path) -> str:
        import pandas as pd
        engine = "openpyxl" if path.suffix.lower() == ".xlsx" else "xlrd"
        try:
            sheets = pd.read_excel(str(path), sheet_name=None, engine=engine)
        except Exception:
            sheets = pd.read_excel(str(path), sheet_name=None)
        parts = []
        for sheet_name, df in sheets.items():
            parts.append(f"[Sheet: {sheet_name}]")
            parts.append(df.to_string(index=False))
        return "\n\n".join(parts)

    @staticmethod
    def _csv(path: Path, ext: str) -> str:
        import pandas as pd
        sep = "\t" if ext == ".tsv" else ","
        df = pd.read_csv(str(path), sep=sep, on_bad_lines="skip")
        return df.to_string(index=False)

    @staticmethod
    def _pptx(path: Path) -> str:
        from pptx import Presentation
        prs = Presentation(str(path))
        parts = []
        for i, slide in enumerate(prs.slides, 1):
            slide_texts = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        t = para.text.strip()
                        if t:
                            slide_texts.append(t)
            if slide_texts:
                parts.append(f"[Slide {i}]\n" + "\n".join(slide_texts))
        return "\n\n".join(parts)

    @staticmethod
    def _html(path: Path) -> str:
        from bs4 import BeautifulSoup
        raw = path.read_bytes()
        soup = BeautifulSoup(raw, "html.parser")
        for tag in soup(["script", "style", "meta", "link"]):
            tag.decompose()
        return soup.get_text(separator="\n", strip=True)

    @staticmethod
    def _eml(path: Path) -> str:
        with open(path, "rb") as f:
            msg = email.message_from_bytes(f.read())
        parts = []
        for key in ("Subject", "From", "To", "Date"):
            val = msg.get(key, "")
            if val:
                parts.append(f"{key}: {val}")
        if msg.is_multipart():
            for part in msg.walk():
                if part.get_content_type() == "text/plain":
                    payload = part.get_payload(decode=True)
                    if payload:
                        parts.append(payload.decode("utf-8", errors="replace"))
        else:
            payload = msg.get_payload(decode=True)
            if payload:
                parts.append(payload.decode("utf-8", errors="replace"))
        return "\n".join(parts)

    @staticmethod
    def _msg(path: Path) -> str:
        """Outlook .msg file (extract-msg)"""
        import extract_msg
        with extract_msg.openMsg(str(path)) as msg:
            parts = []
            if msg.subject:
                parts.append(f"Subject: {msg.subject}")
            if msg.sender:
                parts.append(f"From: {msg.sender}")
            if msg.to:
                parts.append(f"To: {msg.to}")
            if msg.date:
                parts.append(f"Date: {msg.date}")
            if msg.body:
                parts.append(msg.body)
            return "\n".join(parts)

    @staticmethod
    def _json(path: Path) -> str:
        with open(path, encoding="utf-8", errors="replace") as f:
            obj = json.load(f)
        return json.dumps(obj, ensure_ascii=False, indent=2)

    @staticmethod
    def _jsonl(path: Path) -> str:
        lines = []
        with open(path, encoding="utf-8", errors="replace") as f:
            for line in f:
                line = line.strip()
                if line:
                    lines.append(line)
        return "\n".join(lines)

    @staticmethod
    def _rtf(path: Path) -> str:
        """RTF: RTFDE library preferred -> pandoc fallback"""
        try:
            from RTFDE.deencapsulate import DeEncapsulator
            with open(path, "rb") as f:
                raw = f.read()
            dec = DeEncapsulator(raw)
            dec.deencapsulate()
            from bs4 import BeautifulSoup
            return BeautifulSoup(dec.get_content(), "html.parser").get_text(
                separator="\n", strip=True
            )
        except Exception:
            return _TextExtractor._pandoc(path, "plain")

    @staticmethod
    def _odt(path: Path) -> str:
        """LibreOffice ODT (odfpy)"""
        from odf.opendocument import load
        from odf.text import P
        doc = load(str(path))
        parts = []
        for para in doc.getElementsByType(P):
            t = "".join(
                str(node) for node in para.childNodes
                if hasattr(node, "data")
            ).strip()
            if t:
                parts.append(t)
        return "\n".join(parts)

    @staticmethod
    def _archive(path: Path, _depth: int) -> str:
        """
        Extract ZIP / 7z / RAR to a temporary directory and recursively extract text.
        The temporary directory is deleted after extraction.
        """
        ext = path.suffix.lower()
        parts = []
        with tempfile.TemporaryDirectory() as tmp:
            tmp_path = Path(tmp)
            try:
                if ext == ".zip":
                    import zipfile
                    with zipfile.ZipFile(str(path)) as zf:
                        zf.extractall(tmp)
                elif ext == ".7z":
                    import py7zr
                    with py7zr.SevenZipFile(str(path), mode="r") as sz:
                        sz.extractall(path=tmp)
                elif ext == ".rar":
                    import rarfile
                    with rarfile.RarFile(str(path)) as rf:
                        rf.extractall(tmp)
            except Exception as e:
                logger.warning(f"[Extractor] archive extract failed {path.name}: {e}")
                return ""

            for inner in sorted(tmp_path.rglob("*")):
                if inner.is_file():
                    inner_ext = inner.suffix.lower()
                    if inner_ext in _TextExtractor.ARCHIVE_EXTENSIONS:
                        t = _TextExtractor.extract(inner, _depth + 1)
                    else:
                        t = _TextExtractor.extract(inner, _depth)
                    if t.strip():
                        parts.append(f"[{inner.name}]\n{t}")

        return "\n\n".join(parts)

    @staticmethod
    def _pandoc(path: Path, to_fmt: str) -> str:
        try:
            result = subprocess.run(
                ["pandoc", str(path), "-t", to_fmt, "--quiet"],
                capture_output=True, text=True, timeout=30,
            )
            return result.stdout
        except (FileNotFoundError, subprocess.TimeoutExpired) as e:
            logger.warning(f"[Extractor] pandoc failed for {path.name}: {e}")
            return ""


# -- Chunking ----------------------------------------------------------------

_SECTION_SPLIT_RE = re.compile(r'\n(?=# )|^(?=# )|\n---\n', re.MULTILINE)


def _chunk_text(text: str, source_label: str) -> list[dict]:
    """Split text into chunks respecting markdown section boundaries (per §5.2).

    First splits on markdown headers (# ) and horizontal rules (---),
    then applies character-level chunking within each section.
    """
    if not text or len(text.strip()) < MIN_CHUNK_CHARS:
        return []

    # Split into sections at markdown boundaries
    sections = _SECTION_SPLIT_RE.split(text)

    chunks = []
    char_offset = 0
    for section in sections:
        section_flat = " ".join(section.split())
        if len(section_flat) < MIN_CHUNK_CHARS:
            char_offset += len(section)
            continue
        # Sub-chunk within section
        start = 0
        while start < len(section_flat):
            end = min(start + CHUNK_CHARS, len(section_flat))
            chunk_text = section_flat[start:end].strip()
            if len(chunk_text) >= MIN_CHUNK_CHARS:
                chunks.append({
                    "text":         chunk_text,
                    "char_start":   char_offset + start,
                    "char_end":     char_offset + end,
                    "source_label": source_label,
                })
            start += CHUNK_CHARS - OVERLAP_CHARS
            if start >= len(section_flat):
                break
        char_offset += len(section)
    return chunks


# -- Manifest ----------------------------------------------------------------

def _file_hash(path: Path) -> str:
    h = hashlib.sha256()
    with open(path, "rb") as f:
        for block in iter(lambda: f.read(65536), b""):
            h.update(block)
    return h.hexdigest()


def _load_manifest(manifest_path: Path) -> dict:
    if not manifest_path.exists():
        return {"files": {}, "built_at": None}
    try:
        with open(manifest_path, encoding="utf-8") as f:
            return json.load(f)
    except Exception:
        return {"files": {}, "built_at": None}


def _save_manifest(manifest_path: Path, manifest: dict):
    with open(manifest_path, "w", encoding="utf-8") as f:
        json.dump(manifest, f, indent=2, ensure_ascii=False)


# -- IndexBuilder ------------------------------------------------------------

class _IndexBuilder:
    """
    Scans the entire corpus/ folder and builds an IndexFlatIP.
    Uses os.replace() for atomic swap. Embedding device follows
    MMV_EMBEDDING_DEVICE (Principle G).
    """

    SUPPORTED_EXTENSIONS = (
        _TextExtractor.ARCHIVE_EXTENSIONS |
        _TextExtractor.IMAGE_EXTENSIONS |
        {
            ".txt", ".md", ".rst", ".yaml", ".yml",
            ".pdf",
            ".docx", ".doc",
            ".xlsx", ".xls", ".csv", ".tsv",
            ".pptx", ".ppt",
            ".html", ".htm",
            ".eml", ".msg",
            ".json", ".jsonl",
            ".rtf", ".odt",
        }
    )

    def __init__(
        self,
        corpus_dir:    Path,
        index_path:    Path,
        chunks_path:   Path,
        manifest_path: Path,
        model:         SentenceTransformer,
        model_name:    str = MODEL_NAME,
        passage_prefix: str = "",
    ):
        self.corpus_dir    = corpus_dir
        self.index_path    = index_path
        self.chunks_path   = chunks_path
        self.manifest_path = manifest_path
        self.model         = model
        self.model_name    = model_name
        self.passage_prefix = passage_prefix
        self._new_index    = index_path.with_name("custom_index_new.faiss")
        self._new_chunks   = chunks_path.with_name("custom_chunks_new.jsonl")

    def needs_rebuild(self) -> bool:
        if not self.index_path.exists():
            return True
        manifest  = _load_manifest(self.manifest_path)
        old_files = manifest.get("files", {})
        return old_files != self._scan_files()

    def _scan_files(self) -> dict[str, str]:
        result = {}
        if not self.corpus_dir.exists():
            return result
        for path in sorted(self.corpus_dir.rglob("*")):
            if path.is_file() and path.suffix.lower() in self.SUPPORTED_EXTENSIONS:
                rel = str(path.relative_to(self.corpus_dir))
                try:
                    result[rel] = _file_hash(path)
                except Exception:
                    pass
        return result

    def build(self) -> int:
        logger.info("[IndexBuilder] Starting rebuild...")
        t0 = time.time()

        current_files = self._scan_files()
        if not current_files:
            logger.warning(f"[IndexBuilder] corpus_dir empty: {self.corpus_dir}")
            return 0

        all_chunks: list[dict] = []
        all_texts:  list[str]  = []

        for rel_path, file_hash in current_files.items():
            abs_path = self.corpus_dir / rel_path
            logger.info(f"  Extracting: {rel_path}")
            text = _TextExtractor.extract(abs_path)
            if not text.strip():
                logger.warning(f"  [skip] empty: {rel_path}")
                continue
            for ci, chunk in enumerate(_chunk_text(text, rel_path)):
                all_chunks.append({
                    "text":         chunk["text"],
                    "char_start":   chunk["char_start"],
                    "char_end":     chunk["char_end"],
                    "source_path":  str(abs_path),
                    "source_label": rel_path,
                    "chunk_index":  ci,
                    "file_hash":    file_hash,
                })
                all_texts.append(chunk["text"])

        if not all_texts:
            logger.warning("[IndexBuilder] No chunks extracted. Aborting.")
            return 0

        logger.info(f"  {len(all_texts):,} chunks from {len(current_files)} files")

        logger.info(f"  Embedding {len(all_texts):,} chunks...")
        # ME5 convention: prepend "passage: " to corpus chunks when the
        # adapter uses the multilingual-e5 family. No-op for non-ME5
        # experimental overrides.
        if self.passage_prefix:
            texts_to_encode = [self.passage_prefix + t for t in all_texts]
        else:
            texts_to_encode = all_texts
        embeddings = self.model.encode(
            texts_to_encode,
            batch_size           = EMBED_BATCH,
            convert_to_numpy     = True,
            normalize_embeddings = True,
            show_progress_bar    = False,
        ).astype(np.float32)

        # IndexFlatIP (Box A is small, no quantization needed. §4.1)
        index = faiss.IndexFlatIP(embeddings.shape[1])
        index.add(embeddings)

        # Write to temporary files
        faiss.write_index(index, str(self._new_index))
        with open(self._new_chunks, "w", encoding="utf-8") as f:
            for chunk in all_chunks:
                f.write(json.dumps(chunk, ensure_ascii=False) + "\n")

        # Atomic swap (os.replace = assumes same filesystem)
        os.replace(self._new_index,  self.index_path)
        os.replace(self._new_chunks, self.chunks_path)

        # Update manifest
        _save_manifest(self.manifest_path, {
            "files":       current_files,
            "chunk_count": len(all_chunks),
            "built_at":    datetime.now(timezone.utc).isoformat(),
            "model":       self.model_name,
        })

        elapsed = time.time() - t0
        logger.info(f"[IndexBuilder] Done. {len(all_chunks):,} chunks in {elapsed:.1f}s")
        return len(all_chunks)


# -- FileWatcher -------------------------------------------------------------

class _CorpusWatcher:
    """Watches corpus/ with watchdog -> triggers asynchronous rebuild with 3-second debounce"""

    def __init__(self, corpus_dir: Path, on_change):
        self.corpus_dir = corpus_dir
        self.on_change  = on_change
        self._observer  = None
        self._timer: Optional[threading.Timer] = None

    def start(self):
        try:
            from watchdog.observers import Observer
            from watchdog.events import FileSystemEventHandler

            watcher = self

            class _Handler(FileSystemEventHandler):
                def on_any_event(self, event):
                    if event.is_directory:
                        return
                    if Path(event.src_path).suffix.lower() \
                            not in _IndexBuilder.SUPPORTED_EXTENSIONS:
                        return
                    if watcher._timer:
                        watcher._timer.cancel()
                    watcher._timer = threading.Timer(3.0, watcher.on_change)
                    watcher._timer.start()

            self._observer = Observer()
            self._observer.schedule(_Handler(), str(self.corpus_dir), recursive=True)
            self._observer.start()
            logger.info(f"[CorpusWatcher] Watching: {self.corpus_dir}")
        except Exception as e:
            logger.warning(f"[CorpusWatcher] Failed to start: {e}")

    def stop(self):
        if self._timer:
            self._timer.cancel()
        if self._observer:
            self._observer.stop()
            self._observer.join()


# -- CustomRagAdapter --------------------------------------------------------

class CustomRagAdapter(RetrievalAdapter):
    """
    MOBIUS MMV Phase C: Box A Custom RAG adapter.

    Usage:
        adapter = CustomRagAdapter(corpus_dir="corpus/", data_dir="data/box_a/")
        adapter.load()
        result = adapter.retrieve("What is the project deadline?")

    Notes:
        - Call load() once at startup.
        - If corpus/ is empty, is_available()=False (normal behavior. Falls back to Box W).
        - FileWatcher detects corpus/ changes and triggers asynchronous rebuild.
        - During rebuild, continues serving with the old index (Principle F).
        - Embedding device auto-selects through MMV_EMBEDDING_DEVICE.
        - threshold=0.80 (IndexFlatIP cosine; recalibrated for ME5-large).
          ME5 (multilingual-e5-large, "query:"/"passage:" prefixes) has a high
          similarity floor. The legacy MiniLM value (0.50) let every query clear
          the bar, so Box A/0 over-fired and blocked fall-through; 0.80 passes the
          full suite. NOTE: a broad-query probe shows ME5 genuine/unrelated cosine
          *overlap* (genuine ~0.76-0.85, unrelated ~0.72-0.82) — no single
          threshold separates perfectly. The PRIMARY box discrimination is the
          routing layer (config/pattern_library/_taxonomy.yaml: factual_inquiry
          -> box_w, box_0 excluded; self_ref -> box_0). This threshold is a
          secondary sanity net (0.80 = best single value). A margin/contrast
          metric in get_sufficiency_score would be more robust for any path that
          consults Box A/0 outside the router.
    """

    DEFAULT_THRESHOLD = 0.80

    def __init__(
        self,
        corpus_dir: str = "corpus",
        data_dir:   str = "data/box_a",
        threshold:  float = DEFAULT_THRESHOLD,
        watch:      bool = True,
        # Embedding Rule (docs/EMBEDDING_RULE.md): MMV vector indexes use
        # intfloat/multilingual-e5-large. model_name/query_prefix/
        # passage_prefix remain exposed for explicit migration tests or
        # exceptional research runs, but production Box A / Box 0 use ME5.
        model_name:    Optional[str] = None,
        query_prefix:  str = "",
        passage_prefix: str = "",
    ):
        self.corpus_dir     = Path(corpus_dir)
        self.data_dir       = Path(data_dir)
        self.index_path     = self.data_dir / "custom_index.faiss"
        self.chunks_path    = self.data_dir / "custom_chunks.jsonl"
        self.manifest_path  = self.data_dir / "index_manifest.json"
        self.threshold      = threshold
        self._watch         = watch
        self._model_name    = model_name or MODEL_NAME
        _is_me5 = "multilingual-e5" in self._model_name.lower()
        self._query_prefix  = query_prefix or ("query: " if _is_me5 else "")
        self._passage_prefix = passage_prefix or ("passage: " if _is_me5 else "")

        self._model:      Optional[SentenceTransformer] = None
        self._index:      Optional[faiss.Index]          = None
        self._chunks:     list[dict]                     = []
        self._builder:    Optional[_IndexBuilder]        = None
        self._watcher:    Optional[_CorpusWatcher]       = None
        self._lock        = threading.RLock()
        self._loaded      = False
        self._rebuilding  = False
        self._ollama_endpoint = os.environ.get("OLLAMA_ENDPOINT", "http://localhost:11434")
        self._ollama_model = os.environ.get("OLLAMA_MODEL", "qwen3.5:9b")

    def load(self) -> None:
        if self._loaded:
            return
        self.data_dir.mkdir(parents=True, exist_ok=True)
        self.corpus_dir.mkdir(parents=True, exist_ok=True)

        # Detect model-change: if a prior manifest recorded a different
        # embedding model, the on-disk FAISS index has incompatible dim.
        # Force a rebuild by deleting the stale index artifacts before the
        # IndexBuilder is invoked. This preserves the corpus/ directory
        # itself; only derived artifacts are removed.
        if self.manifest_path.exists():
            try:
                _prior = _load_manifest(self.manifest_path)
                _prior_model = _prior.get("model")
                if _prior_model and _prior_model != self._model_name:
                    logger.info(
                        f"[CustomRagAdapter] Embedding model change detected "
                        f"({_prior_model} -> {self._model_name}); "
                        f"clearing stale index artifacts."
                    )
                    for _p in (self.index_path, self.chunks_path, self.manifest_path):
                        if _p.exists():
                            _p.unlink()
            except Exception as _e:
                logger.warning(
                    f"[CustomRagAdapter] Could not read prior manifest: {_e}"
                )

        device = _select_embedding_device()
        logger.info(f"[CustomRagAdapter] Loading model: {self._model_name} (device={device})")
        try:
            self._model = SentenceTransformer(self._model_name, device=device)
        except RuntimeError as e:
            # CUDA OOM / device contention (e.g. Ollama holding the GPU):
            # degrade to CPU rather than crash the runtime.
            if device != "cpu":
                logger.warning(
                    f"[CustomRagAdapter] model load on {device} failed "
                    f"({type(e).__name__}: {e}); falling back to CPU."
                )
                self._model = SentenceTransformer(self._model_name, device="cpu")
            else:
                raise

        self._builder = _IndexBuilder(
            corpus_dir     = self.corpus_dir,
            index_path     = self.index_path,
            chunks_path    = self.chunks_path,
            manifest_path  = self.manifest_path,
            model          = self._model,
            model_name     = self._model_name,
            passage_prefix = self._passage_prefix,
        )

        if self._builder.needs_rebuild():
            logger.info("[CustomRagAdapter] Building index...")
            self._builder.build()

        self._load_index()

        if self._watch:
            self._watcher = _CorpusWatcher(self.corpus_dir, self._trigger_rebuild)
            self._watcher.start()

        self._loaded = True
        logger.info(
            f"[CustomRagAdapter] Ready. "
            f"{len(self._chunks):,} chunks, threshold={self.threshold}"
        )

    @staticmethod
    def _looks_english(text: str) -> bool:
        try:
            from langdetect import detect
            return detect(text) == "en"
        except Exception:
            ascii_count = sum(1 for c in text if ord(c) < 128)
            return len(text) > 0 and ascii_count / len(text) > 0.95

    def _translate_to_english(self, query: str) -> str:
        """Translate non-English query to English via Ollama. Fail-safe."""
        if self._looks_english(query):
            return query
        try:
            import requests as _req
            resp = _req.post(
                f"{self._ollama_endpoint}/api/generate",
                json={
                    "model": self._ollama_model,
                    "prompt": (
                        "Translate the following to English. "
                        "Output ONLY the translation, nothing else.\n\n"
                        f"{query}"
                    ),
                    "stream": False,
                    "think": False,
                    "options": {"temperature": 0.0, "num_predict": 64},
                },
                timeout=15,
            )
            resp.raise_for_status()
            translated = resp.json().get("response", "").strip()
            if translated:
                return translated
        except Exception:
            pass
        return query

    def retrieve(self, query: str, top_k: int = DEFAULT_TOP_K) -> RetrievalResult:
        if not self._loaded:
            raise RuntimeError("CustomRagAdapter not loaded. Call load() first.")

        top_k = max(1, min(top_k, MAX_TOP_K))
        now   = datetime.now(timezone.utc).isoformat()

        with self._lock:
            index  = self._index
            chunks = self._chunks

        if index is None or not chunks:
            return RetrievalResult(sources=[], outcome="failed", synthesis="")

        # ME5 is cross-lingual by design (JA/EN/ZH share one vector space), so
        # we skip the Ollama English-translation hop when an ME5 query prefix
        # is configured. For non-ME5 experimental overrides (no prefix),
        # translation is still attempted.
        if self._query_prefix:
            query_for_search = self._query_prefix + query
        else:
            query_for_search = self._translate_to_english(query)

        query_vec = self._model.encode(
            [query_for_search],
            convert_to_numpy     = True,
            normalize_embeddings = True,
            show_progress_bar    = False,
        ).astype(np.float32)

        actual_k = min(top_k, index.ntotal)
        scores_arr, indices_arr = index.search(query_vec, actual_k)

        sources: list[Source] = []
        texts:   list[str]    = []

        for i, idx in enumerate(indices_arr[0]):
            if idx < 0 or idx >= len(chunks):
                continue
            score = float(np.clip(scores_arr[0][i], 0.0, 1.0))
            chunk = chunks[idx]
            sources.append(Source(
                source_type     = "local_rag",
                label           = chunk.get("source_label", ""),
                uri             = "file://" + chunk.get("source_path", ""),
                chunk_index     = chunk.get("chunk_index"),
                retrieved_at    = now,
                relevance_score = round(score, 4),
            ))
            texts.append(chunk["text"])

        if not sources:
            return RetrievalResult(sources=[], outcome="failed", synthesis="")

        outcome = "success" if len(sources) == top_k else "partial"
        return RetrievalResult(
            sources   = sources,
            outcome   = outcome,
            synthesis = "\n\n".join(texts),
        )

    def is_available(self) -> bool:
        if not self._loaded:
            return False
        with self._lock:
            return self._index is not None and len(self._chunks) > 0

    def get_sufficiency_score(self, result: RetrievalResult) -> float:
        if not result.sources:
            return 0.0
        scores = [s.relevance_score for s in result.sources if s.relevance_score is not None]
        return max(scores) if scores else 0.0

    def rebuild(self) -> int:
        return self._trigger_rebuild(blocking=True)

    def close(self) -> None:
        if self._watcher:
            self._watcher.stop()
        with self._lock:
            self._index  = None
            self._chunks = []
        self._model  = None
        self._loaded = False

    def _load_index(self):
        if not self.index_path.exists():
            logger.info("[CustomRagAdapter] No index. corpus/ may be empty.")
            return
        try:
            new_index  = faiss.read_index(str(self.index_path))
            new_chunks = []
            with open(self.chunks_path, encoding="utf-8") as f:
                for line in f:
                    line = line.strip()
                    if line:
                        new_chunks.append(json.loads(line))
            with self._lock:
                self._index  = new_index
                self._chunks = new_chunks
            logger.info(
                f"[CustomRagAdapter] Index loaded: "
                f"{new_index.ntotal:,} vectors, {len(new_chunks):,} chunks"
            )
        except Exception as e:
            logger.error(f"[CustomRagAdapter] Index load failed: {e}")

    def _trigger_rebuild(self, blocking: bool = False) -> int:
        if self._rebuilding:
            logger.info("[CustomRagAdapter] Rebuild already in progress.")
            return 0

        def _run():
            self._rebuilding = True
            try:
                n = self._builder.build()
                if n > 0:
                    self._load_index()
                    logger.info(f"[CustomRagAdapter] Rebuild complete. {n:,} chunks.")
                return n
            except Exception as e:
                logger.error(f"[CustomRagAdapter] Rebuild failed: {e}")
                return 0
            finally:
                self._rebuilding = False

        if blocking:
            return _run()
        t = threading.Thread(target=_run, daemon=True, name="CustomRagRebuild")
        t.start()
        return 0

    def __repr__(self) -> str:
        status = "loaded" if self._loaded else "not loaded"
        with self._lock:
            n = len(self._chunks)
        return (
            f"CustomRagAdapter({status}, "
            f"corpus={self.corpus_dir}, "
            f"chunks={n:,}, threshold={self.threshold})"
        )


# -- CLI ---------------------------------------------------------------------

def _cli():
    import argparse
    parser = argparse.ArgumentParser(
        description="CustomRagAdapter CLI — Box A search test / index building"
    )
    parser.add_argument("--corpus",    default="corpus")
    parser.add_argument("--data-dir",  default="data/box_a")
    parser.add_argument("--top-k",     type=int, default=5)
    parser.add_argument("--threshold", type=float, default=CustomRagAdapter.DEFAULT_THRESHOLD)
    parser.add_argument("--build",     action="store_true", help="Force rebuild and exit")
    parser.add_argument("query", nargs="?", default=None)
    args = parser.parse_args()

    logging.basicConfig(level=logging.INFO, format="%(asctime)s %(levelname)s %(message)s")

    adapter = CustomRagAdapter(
        corpus_dir=args.corpus, data_dir=args.data_dir,
        threshold=args.threshold, watch=False,
    )
    adapter.load()

    if args.build:
        print(f"Built: {adapter.rebuild():,} chunks")
        return

    if not adapter.is_available():
        print(f"[warn] No files found in corpus/: {args.corpus}")
        return

    query = args.query or "test"
    print(f"\nQuery: {query}\nAdapter: {adapter}")
    t0 = time.time()
    result = adapter.retrieve(query, top_k=args.top_k)
    elapsed = (time.time() - t0) * 1000
    score = adapter.get_sufficiency_score(result)
    print(f"\nOutcome: {result.outcome}  Score: {score:.4f}  Latency: {elapsed:.1f}ms")
    blocks = result.synthesis.split("\n\n")
    for i, src in enumerate(result.sources, 1):
        snippet = blocks[i - 1][:120] if i - 1 < len(blocks) else ""
        print(f"  [{i}] {src.relevance_score:.4f}  {src.label}\n       {snippet}...")


if __name__ == "__main__":
    _cli()
